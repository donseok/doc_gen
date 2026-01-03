"""
템플릿 분석 서비스

업로드된 PPTX 템플릿을 분석하여 디자인 요소를 추출합니다.
- 색상 팔레트
- 폰트 설정
- 레이아웃 패턴
- 슬라이드 구조
"""
from typing import Dict, List, Optional, Any
from dataclasses import dataclass, field, asdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os


@dataclass
class ColorPalette:
    """색상 팔레트"""
    primary: str = "#1E3A8A"       # 주 색상
    secondary: str = "#64748B"     # 보조 색상
    accent: str = "#3B82F6"        # 강조 색상
    background: str = "#FFFFFF"    # 배경 색상
    text_dark: str = "#1E293B"     # 어두운 텍스트
    text_light: str = "#FFFFFF"    # 밝은 텍스트
    additional_colors: List[str] = field(default_factory=list)


@dataclass
class FontSettings:
    """폰트 설정"""
    title_font: str = "맑은 고딕"
    body_font: str = "맑은 고딕"
    english_font: str = "Arial"
    title_size: int = 44           # pt
    subtitle_size: int = 28
    body_size: int = 18
    small_size: int = 14


@dataclass
class LayoutInfo:
    """레이아웃 정보"""
    name: str
    index: int
    has_title: bool = False
    has_content: bool = False
    has_image: bool = False
    placeholders: List[Dict] = field(default_factory=list)


@dataclass
class TemplateStyle:
    """템플릿 스타일 정보"""
    name: str
    colors: ColorPalette
    fonts: FontSettings
    layouts: List[LayoutInfo]
    slide_width: float = 13.333    # inches
    slide_height: float = 7.5      # inches
    margin_left: float = 0.5
    margin_right: float = 0.5
    margin_top: float = 0.5
    margin_bottom: float = 0.5
    header_height: float = 0.15    # 상단 바 높이
    use_header_bar: bool = True
    use_footer: bool = False

    def to_dict(self) -> Dict:
        """딕셔너리로 변환"""
        return {
            "name": self.name,
            "colors": asdict(self.colors),
            "fonts": asdict(self.fonts),
            "layouts": [asdict(l) for l in self.layouts],
            "slide_width": self.slide_width,
            "slide_height": self.slide_height,
            "margin_left": self.margin_left,
            "margin_right": self.margin_right,
            "margin_top": self.margin_top,
            "margin_bottom": self.margin_bottom,
            "header_height": self.header_height,
            "use_header_bar": self.use_header_bar,
            "use_footer": self.use_footer,
        }


class TemplateAnalyzer:
    """템플릿 분석기"""

    def __init__(self):
        self.extracted_colors: List[str] = []
        self.extracted_fonts: List[str] = []

    def analyze(self, pptx_path: str) -> TemplateStyle:
        """
        PPTX 파일을 분석하여 스타일 정보 추출

        Args:
            pptx_path: PPTX 파일 경로

        Returns:
            TemplateStyle: 추출된 스타일 정보
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"템플릿 파일을 찾을 수 없습니다: {pptx_path}")

        prs = Presentation(pptx_path)

        # 슬라이드 크기 추출
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches

        # 레이아웃 분석
        layouts = self._analyze_layouts(prs)

        # 슬라이드 내용 분석 (색상, 폰트)
        self._analyze_slides(prs)

        # 색상 팔레트 생성
        colors = self._create_color_palette()

        # 폰트 설정 생성
        fonts = self._create_font_settings()

        # 마진 및 헤더 분석
        margins = self._analyze_margins(prs)

        template_name = os.path.splitext(os.path.basename(pptx_path))[0]

        return TemplateStyle(
            name=template_name,
            colors=colors,
            fonts=fonts,
            layouts=layouts,
            slide_width=slide_width,
            slide_height=slide_height,
            margin_left=margins.get("left", 0.5),
            margin_right=margins.get("right", 0.5),
            margin_top=margins.get("top", 0.5),
            margin_bottom=margins.get("bottom", 0.5),
            header_height=margins.get("header_height", 0.15),
            use_header_bar=margins.get("use_header_bar", True),
            use_footer=margins.get("use_footer", False),
        )

    def _analyze_layouts(self, prs: Presentation) -> List[LayoutInfo]:
        """레이아웃 정보 분석"""
        layouts = []

        for idx, layout in enumerate(prs.slide_layouts):
            placeholders = []
            has_title = False
            has_content = False
            has_image = False

            for shape in layout.placeholders:
                ph_info = {
                    "idx": shape.placeholder_format.idx,
                    "type": str(shape.placeholder_format.type),
                    "left": shape.left.inches if shape.left else 0,
                    "top": shape.top.inches if shape.top else 0,
                    "width": shape.width.inches if shape.width else 0,
                    "height": shape.height.inches if shape.height else 0,
                }
                placeholders.append(ph_info)

                # 플레이스홀더 타입 확인
                ph_type = str(shape.placeholder_format.type)
                if "TITLE" in ph_type or "CENTER_TITLE" in ph_type:
                    has_title = True
                elif "BODY" in ph_type or "OBJECT" in ph_type:
                    has_content = True
                elif "PICTURE" in ph_type:
                    has_image = True

            layouts.append(LayoutInfo(
                name=layout.name,
                index=idx,
                has_title=has_title,
                has_content=has_content,
                has_image=has_image,
                placeholders=placeholders,
            ))

        return layouts

    def _analyze_slides(self, prs: Presentation):
        """슬라이드 내용 분석 (색상, 폰트 추출)"""
        for slide in prs.slides:
            for shape in slide.shapes:
                self._extract_shape_styles(shape)

    def _extract_shape_styles(self, shape):
        """도형에서 스타일 추출"""
        # 채우기 색상 추출
        if hasattr(shape, "fill") and shape.fill:
            try:
                if shape.fill.type is not None:
                    fore_color = shape.fill.fore_color
                    if fore_color and fore_color.type is not None:
                        if hasattr(fore_color, "rgb") and fore_color.rgb:
                            color_hex = f"#{fore_color.rgb}"
                            if color_hex not in self.extracted_colors:
                                self.extracted_colors.append(color_hex)
            except Exception:
                pass

        # 테두리 색상 추출
        if hasattr(shape, "line") and shape.line:
            try:
                if shape.line.color and shape.line.color.type is not None:
                    if hasattr(shape.line.color, "rgb") and shape.line.color.rgb:
                        color_hex = f"#{shape.line.color.rgb}"
                        if color_hex not in self.extracted_colors:
                            self.extracted_colors.append(color_hex)
            except Exception:
                pass

        # 텍스트 프레임에서 폰트 및 색상 추출
        if hasattr(shape, "text_frame"):
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 폰트 추출
                        if run.font.name and run.font.name not in self.extracted_fonts:
                            self.extracted_fonts.append(run.font.name)

                        # 텍스트 색상 추출
                        if run.font.color and run.font.color.type is not None:
                            if hasattr(run.font.color, "rgb") and run.font.color.rgb:
                                color_hex = f"#{run.font.color.rgb}"
                                if color_hex not in self.extracted_colors:
                                    self.extracted_colors.append(color_hex)
            except Exception:
                pass

    def _create_color_palette(self) -> ColorPalette:
        """추출된 색상으로 팔레트 생성"""
        palette = ColorPalette()

        if not self.extracted_colors:
            return palette

        # 색상 분류 (밝기 기준)
        dark_colors = []
        light_colors = []
        mid_colors = []

        for color in self.extracted_colors:
            brightness = self._get_color_brightness(color)
            if brightness < 85:
                dark_colors.append(color)
            elif brightness > 200:
                light_colors.append(color)
            else:
                mid_colors.append(color)

        # 팔레트 구성
        if dark_colors:
            palette.primary = dark_colors[0]
            palette.text_dark = dark_colors[0]
            if len(dark_colors) > 1:
                palette.secondary = dark_colors[1]

        if mid_colors:
            palette.accent = mid_colors[0]

        if light_colors:
            palette.background = light_colors[0] if light_colors[0] != "#FFFFFF" else palette.background

        palette.additional_colors = self.extracted_colors[:10]  # 최대 10개

        return palette

    def _get_color_brightness(self, hex_color: str) -> float:
        """색상의 밝기 계산 (0-255)"""
        try:
            hex_color = hex_color.lstrip("#")
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return (r * 299 + g * 587 + b * 114) / 1000
        except Exception:
            return 128  # 중간값 반환

    def _create_font_settings(self) -> FontSettings:
        """추출된 폰트로 설정 생성"""
        settings = FontSettings()

        if not self.extracted_fonts:
            return settings

        # 한글 폰트와 영문 폰트 분류
        korean_fonts = []
        english_fonts = []

        for font in self.extracted_fonts:
            if self._is_korean_font(font):
                korean_fonts.append(font)
            else:
                english_fonts.append(font)

        # 폰트 설정
        if korean_fonts:
            settings.title_font = korean_fonts[0]
            settings.body_font = korean_fonts[0]

        if english_fonts:
            settings.english_font = english_fonts[0]

        return settings

    def _is_korean_font(self, font_name: str) -> bool:
        """한글 폰트인지 확인"""
        korean_indicators = [
            "맑은", "고딕", "굴림", "돋움", "바탕", "나눔",
            "Malgun", "Gulim", "Dotum", "Batang", "Nanum"
        ]
        return any(indicator in font_name for indicator in korean_indicators)

    def _analyze_margins(self, prs: Presentation) -> Dict[str, float]:
        """슬라이드 마진 분석"""
        margins = {
            "left": 0.5,
            "right": 0.5,
            "top": 0.5,
            "bottom": 0.5,
            "header_height": 0.15,
            "use_header_bar": False,
            "use_footer": False,
        }

        if not prs.slides:
            return margins

        # 첫 번째 슬라이드에서 마진 추정
        slide = prs.slides[0]

        min_left = prs.slide_width.inches
        max_right = 0
        min_top = prs.slide_height.inches
        max_bottom = 0

        for shape in slide.shapes:
            if hasattr(shape, "left") and shape.left:
                left = shape.left.inches
                right = left + (shape.width.inches if shape.width else 0)
                top = shape.top.inches if shape.top else 0
                bottom = top + (shape.height.inches if shape.height else 0)

                if left < min_left:
                    min_left = left
                if right > max_right:
                    max_right = right
                if top < min_top:
                    min_top = top
                if bottom > max_bottom:
                    max_bottom = bottom

                # 상단 얇은 바 감지 (헤더 바)
                if top < 0.3 and (shape.height.inches if shape.height else 0) < 0.3:
                    margins["use_header_bar"] = True
                    margins["header_height"] = shape.height.inches if shape.height else 0.15

        # 마진 계산
        if min_left < prs.slide_width.inches:
            margins["left"] = max(0.25, min_left)
        if max_right > 0:
            margins["right"] = max(0.25, prs.slide_width.inches - max_right)
        if min_top < prs.slide_height.inches:
            margins["top"] = max(0.25, min_top)
        if max_bottom > 0:
            margins["bottom"] = max(0.25, prs.slide_height.inches - max_bottom)

        return margins


def analyze_template(pptx_path: str) -> Dict[str, Any]:
    """
    템플릿 분석 헬퍼 함수

    Args:
        pptx_path: PPTX 파일 경로

    Returns:
        Dict: 템플릿 스타일 정보 딕셔너리
    """
    analyzer = TemplateAnalyzer()
    style = analyzer.analyze(pptx_path)
    return style.to_dict()
