"""
PPT 생성 서비스

개선 사항 (제미나이 분석서 반영):
- 레이아웃 인덱스 기반 → 레이아웃 이름 기반 매핑으로 전환
- 커스텀 템플릿 지원을 위한 유연한 레이아웃 시스템
- 템플릿별 메타데이터 정의 가능
"""
import os
import uuid
from typing import Dict, List, Any, Optional
from datetime import datetime
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.config import settings


@dataclass
class LayoutMapping:
    """
    슬라이드 레이아웃 매핑 정보

    인덱스 기반이 아닌 이름 기반으로 레이아웃을 찾아 적용
    """
    title: str = "Title Slide"
    title_content: str = "Title and Content"
    section_header: str = "Section Header"
    two_content: str = "Two Content"
    comparison: str = "Comparison"
    title_only: str = "Title Only"
    blank: str = "Blank"
    content_caption: str = "Content with Caption"

    # 한글 레이아웃 이름 대체 (한글 템플릿 지원)
    alternatives: Dict[str, List[str]] = field(default_factory=lambda: {
        "title": ["Title Slide", "제목 슬라이드", "표지"],
        "title_content": ["Title and Content", "제목 및 내용", "제목과 내용"],
        "section_header": ["Section Header", "구역 머리글", "섹션 머리글"],
        "two_content": ["Two Content", "두 개의 콘텐츠", "2단 콘텐츠"],
        "comparison": ["Comparison", "비교", "비교형"],
        "title_only": ["Title Only", "제목만", "제목 전용"],
        "blank": ["Blank", "빈 화면", "빈 슬라이드"],
        "content_caption": ["Content with Caption", "캡션 있는 콘텐츠"],
    })


class LayoutResolver:
    """
    템플릿에서 레이아웃을 이름으로 검색하는 리졸버

    인덱스 하드코딩 문제 해결:
    - 레이아웃 이름으로 검색 후 없으면 대체 이름 시도
    - 모든 대체 이름 실패 시 인덱스 폴백
    """

    def __init__(self, presentation: Presentation, mapping: Optional[LayoutMapping] = None):
        self.prs = presentation
        self.mapping = mapping or LayoutMapping()
        self._layout_cache: Dict[str, Any] = {}
        self._build_cache()

    def _build_cache(self):
        """레이아웃 이름 → 레이아웃 객체 캐시 구축"""
        for layout in self.prs.slide_layouts:
            name = layout.name.strip()
            self._layout_cache[name.lower()] = layout

    def get_layout(self, layout_type: str) -> Any:
        """
        레이아웃 타입으로 적절한 레이아웃 객체 반환

        Args:
            layout_type: 레이아웃 유형 (title, title_content, blank 등)

        Returns:
            SlideLayout 객체
        """
        # 대체 이름 목록 가져오기
        alternatives = self.mapping.alternatives.get(layout_type, [])

        # 각 대체 이름으로 검색 시도
        for name in alternatives:
            layout = self._layout_cache.get(name.lower())
            if layout is not None:
                return layout

        # 폴백: 인덱스 기반 (기본 PowerPoint 템플릿 호환)
        fallback_indices = {
            "title": 0,
            "title_content": 1,
            "section_header": 2,
            "two_content": 3,
            "comparison": 4,
            "title_only": 5,
            "blank": 6,
            "content_caption": 7,
        }

        idx = fallback_indices.get(layout_type, 1)
        try:
            return self.prs.slide_layouts[idx]
        except IndexError:
            # 최소한 첫 번째 레이아웃이라도 반환
            return self.prs.slide_layouts[0]

    def list_available_layouts(self) -> List[Dict[str, Any]]:
        """사용 가능한 레이아웃 목록 반환 (디버깅/문서화용)"""
        return [
            {"index": i, "name": layout.name}
            for i, layout in enumerate(self.prs.slide_layouts)
        ]


class PPTGenerator:
    """
    파싱된 마크다운 데이터를 PPT 파일로 생성

    개선 사항:
    - 레이아웃 이름 기반 매핑으로 커스텀 템플릿 호환성 향상
    - LayoutResolver를 통한 유연한 레이아웃 검색
    """

    def __init__(self, layout_mapping: Optional[LayoutMapping] = None):
        self.layout_mapping = layout_mapping or LayoutMapping()
        self.default_font = "맑은 고딕"
        self.title_font_size = Pt(44)
        self.content_font_size = Pt(24)
        self.bullet_font_size = Pt(20)

    def generate(
        self,
        parsed_data: Dict[str, Any],
        template_id: Optional[int] = None,
        options: Optional[Dict[str, Any]] = None,
    ) -> str:
        """
        파싱된 데이터로 PPT 생성

        Args:
            parsed_data: 마크다운 파서에서 생성된 슬라이드 데이터
            template_id: 사용할 템플릿 ID (옵션)
            options: 생성 옵션

        Returns:
            생성된 PPT 파일 경로
        """
        options = options or {}

        # 프레젠테이션 생성
        if template_id:
            prs = self._load_template(template_id)
        else:
            prs = Presentation()

        # 레이아웃 리졸버 초기화
        layout_resolver = LayoutResolver(prs, self.layout_mapping)

        # 제목 슬라이드 추가
        if parsed_data.get("title"):
            self._add_title_slide(prs, parsed_data["title"], layout_resolver)

        # 컨텐츠 슬라이드 추가
        for slide_data in parsed_data.get("slides", []):
            self._add_content_slide(prs, slide_data, layout_resolver)

        # 파일 저장
        output_path = self._generate_output_path(parsed_data.get("title", "presentation"))
        prs.save(output_path)

        return output_path

    def _load_template(self, template_id: int) -> Presentation:
        """템플릿 파일 로드"""
        # TODO: 데이터베이스에서 템플릿 경로 조회
        template_path = os.path.join(settings.TEMPLATES_DIR, "default", "template.pptx")

        if os.path.exists(template_path):
            return Presentation(template_path)

        return Presentation()

    def _add_title_slide(self, prs: Presentation, title: str, resolver: LayoutResolver):
        """제목 슬라이드 추가"""
        slide_layout = resolver.get_layout("title")
        slide = prs.slides.add_slide(slide_layout)

        # 제목 설정
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._format_title(slide.shapes.title)

        # 부제목 (날짜)
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # 부제목 placeholder
                shape.text = datetime.now().strftime("%Y년 %m월 %d일")
                break

    def _add_content_slide(self, prs: Presentation, slide_data: Dict[str, Any], resolver: LayoutResolver):
        """컨텐츠 슬라이드 추가"""
        layout_type = slide_data.get("layout", "title_content")

        # 레이아웃 타입 매핑
        layout_map = {
            "image_only": "blank",
            "image_content": "two_content",
            "code": "title_content",
            "title_content": "title_content",
        }
        resolved_layout = layout_map.get(layout_type, "title_content")
        slide_layout = resolver.get_layout(resolved_layout)

        slide = prs.slides.add_slide(slide_layout)

        # 제목 설정
        if slide.shapes.title and slide_data.get("title"):
            slide.shapes.title.text = slide_data["title"]
            self._format_title(slide.shapes.title)

        # 컨텐츠 추가
        content = slide_data.get("content", [])
        if content:
            self._add_content_to_slide(slide, content)

        # 이미지 추가
        images = slide_data.get("images", [])
        for img in images:
            self._add_image_to_slide(slide, img)

        # 코드 블록 추가
        code_blocks = slide_data.get("code_blocks", [])
        for code_block in code_blocks:
            self._add_code_block_to_slide(slide, code_block)

    def _add_content_to_slide(self, slide, content: List[Dict[str, Any]]):
        """슬라이드에 컨텐츠 추가"""
        # 컨텐츠 placeholder 찾기
        content_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # 컨텐츠 placeholder
                content_shape = shape
                break

        if not content_shape:
            # Placeholder가 없으면 텍스트 박스 생성
            content_shape = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(9), Inches(5)
            )

        tf = content_shape.text_frame
        tf.clear()

        first_paragraph = True
        for item in content:
            if first_paragraph:
                p = tf.paragraphs[0]
                first_paragraph = False
            else:
                p = tf.add_paragraph()

            p.text = item.get("text", "")

            # 항목 타입에 따른 포맷팅
            item_type = item.get("type", "text")
            if item_type == "bullet":
                p.level = 0
                p.font.size = self.bullet_font_size
            elif item_type == "numbered":
                p.level = 0
                p.font.size = self.bullet_font_size
            elif item_type == "subheading":
                p.font.bold = True
                p.font.size = Pt(28)
            elif item_type == "quote":
                p.font.italic = True
                p.font.size = self.content_font_size
            else:
                p.font.size = self.content_font_size

            p.font.name = self.default_font

    def _add_image_to_slide(self, slide, image_data: Dict[str, str]):
        """슬라이드에 이미지 추가"""
        src = image_data.get("src", "")

        # 로컬 파일 경로인 경우
        if os.path.exists(src):
            try:
                slide.shapes.add_picture(
                    src,
                    Inches(5), Inches(2),
                    width=Inches(4)
                )
            except Exception as e:
                print(f"이미지 추가 실패: {e}")

    def _add_code_block_to_slide(self, slide, code_block: Dict[str, str]):
        """슬라이드에 코드 블록 추가"""
        code = code_block.get("code", "")

        # 코드 박스 생성
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(4)

        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = code
        p.font.name = "Consolas"
        p.font.size = Pt(14)

        # 배경색 설정 (라이트 그레이)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)

    def _format_title(self, title_shape):
        """제목 포맷팅"""
        if hasattr(title_shape, 'text_frame'):
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.name = self.default_font
                paragraph.font.size = self.title_font_size
                paragraph.font.bold = True

    def _generate_output_path(self, title: str) -> str:
        """출력 파일 경로 생성"""
        # 안전한 파일명 생성
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_title = safe_title[:50] if safe_title else "presentation"

        # 고유 ID 추가
        unique_id = str(uuid.uuid4())[:8]
        filename = f"{safe_title}_{unique_id}.pptx"

        output_path = os.path.join(settings.OUTPUT_DIR, filename)

        # 출력 디렉토리 생성
        os.makedirs(settings.OUTPUT_DIR, exist_ok=True)

        return output_path
