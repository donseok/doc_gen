"""
PPT 디자인 에이전트

요약된 콘텐츠를 시각적으로 깔끔한 PPT로 변환하는 에이전트
- 모던한 색상 테마 적용
- 업로드된 템플릿 스타일 적용 지원
- 레이아웃별 최적화된 디자인
- 텍스트 오버플로우 방지
- 아이콘/도형 활용
"""
from typing import Dict, List, Any, Optional
from dataclasses import dataclass
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

from app.services.content_summarizer import PresentationContent, SlideContent


@dataclass
class ColorTheme:
    """색상 테마"""
    primary: RGBColor      # 주 색상 (제목, 강조)
    secondary: RGBColor    # 보조 색상
    accent: RGBColor       # 액센트 색상
    background: RGBColor   # 배경색
    text_dark: RGBColor    # 어두운 텍스트
    text_light: RGBColor   # 밝은 텍스트
    gradient_start: RGBColor  # 그라데이션 시작
    gradient_end: RGBColor    # 그라데이션 끝


class PPTDesigner:
    """
    PPT 디자인 에이전트

    요약된 콘텐츠를 시각적으로 매력적인 PPT로 변환
    """

    # 슬라이드 크기 (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # 여백
    MARGIN_LEFT = Inches(0.7)
    MARGIN_RIGHT = Inches(0.7)
    MARGIN_TOP = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)

    # 콘텐츠 영역
    CONTENT_WIDTH = Inches(12)
    CONTENT_TOP = Inches(1.8)
    CONTENT_HEIGHT = Inches(5.2)

    # 폰트
    FONT_TITLE = "맑은 고딕"
    FONT_BODY = "맑은 고딕"
    FONT_ENGLISH = "Segoe UI"

    # 폰트 크기
    FONT_SIZE_MAIN_TITLE = Pt(44)
    FONT_SIZE_SLIDE_TITLE = Pt(32)
    FONT_SIZE_SUBTITLE = Pt(20)
    FONT_SIZE_BODY = Pt(18)
    FONT_SIZE_BULLET = Pt(16)
    FONT_SIZE_SMALL = Pt(14)
    FONT_SIZE_METRIC_VALUE = Pt(36)
    FONT_SIZE_METRIC_LABEL = Pt(14)

    # 불릿당 최대 줄 수
    MAX_LINES_PER_BULLET = 2
    # 줄당 최대 문자 수
    MAX_CHARS_PER_LINE = 35

    def __init__(self, theme: str = "modern_blue", template_style: Optional[Dict] = None):
        """
        Args:
            theme: 색상 테마 (modern_blue, corporate, dark, minimal)
            template_style: 업로드된 템플릿에서 추출한 스타일 정보 (선택적)
        """
        if template_style:
            self.theme = self._apply_template_style(template_style)
            self._apply_template_settings(template_style)
        else:
            self.theme = self._get_theme(theme)
        self.prs = None
        self.template_style = template_style

    def _apply_template_style(self, style: Dict) -> ColorTheme:
        """템플릿 스타일에서 ColorTheme 생성"""
        colors = style.get("colors", {})

        def hex_to_rgb(hex_color: str) -> RGBColor:
            """HEX 색상을 RGBColor로 변환"""
            hex_color = hex_color.lstrip("#")
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)

        return ColorTheme(
            primary=hex_to_rgb(colors.get("primary", "#0070C0")),
            secondary=hex_to_rgb(colors.get("secondary", "#44546A")),
            accent=hex_to_rgb(colors.get("accent", "#00B0F0")),
            background=hex_to_rgb(colors.get("background", "#FFFFFF")),
            text_dark=hex_to_rgb(colors.get("text_dark", "#333333")),
            text_light=hex_to_rgb(colors.get("text_light", "#FFFFFF")),
            gradient_start=hex_to_rgb(colors.get("primary", "#0070C0")),
            gradient_end=hex_to_rgb(colors.get("accent", "#00B0F0")),
        )

    def _apply_template_settings(self, style: Dict):
        """템플릿 설정 적용 (폰트, 마진 등)"""
        fonts = style.get("fonts", {})

        # 폰트 설정 적용
        if fonts.get("title_font"):
            self.FONT_TITLE = fonts["title_font"]
        if fonts.get("body_font"):
            self.FONT_BODY = fonts["body_font"]
        if fonts.get("english_font"):
            self.FONT_ENGLISH = fonts["english_font"]

        # 폰트 크기 적용
        if fonts.get("title_size"):
            self.FONT_SIZE_MAIN_TITLE = Pt(fonts["title_size"])
        if fonts.get("body_size"):
            self.FONT_SIZE_BODY = Pt(fonts["body_size"])
            self.FONT_SIZE_BULLET = Pt(fonts["body_size"] - 2)

        # 마진 적용
        if style.get("margin_left"):
            self.MARGIN_LEFT = Inches(style["margin_left"])
        if style.get("margin_right"):
            self.MARGIN_RIGHT = Inches(style["margin_right"])
        if style.get("margin_top"):
            self.MARGIN_TOP = Inches(style["margin_top"])

        # 헤더 바 설정
        self.use_header_bar = style.get("use_header_bar", True)
        self.header_height = Inches(style.get("header_height", 0.15))

    def _get_theme(self, theme_name: str) -> ColorTheme:
        """테마 색상 반환"""
        themes = {
            "modern_blue": ColorTheme(
                primary=RGBColor(0, 112, 192),      # 파란색
                secondary=RGBColor(68, 84, 106),    # 진한 회색
                accent=RGBColor(0, 176, 240),       # 밝은 파란색
                background=RGBColor(255, 255, 255), # 흰색
                text_dark=RGBColor(51, 51, 51),     # 진한 회색
                text_light=RGBColor(255, 255, 255), # 흰색
                gradient_start=RGBColor(0, 112, 192),
                gradient_end=RGBColor(0, 176, 240),
            ),
            "corporate": ColorTheme(
                primary=RGBColor(0, 51, 102),       # 네이비
                secondary=RGBColor(128, 128, 128),  # 회색
                accent=RGBColor(255, 153, 0),       # 주황색
                background=RGBColor(255, 255, 255),
                text_dark=RGBColor(51, 51, 51),
                text_light=RGBColor(255, 255, 255),
                gradient_start=RGBColor(0, 51, 102),
                gradient_end=RGBColor(0, 102, 153),
            ),
            "dark": ColorTheme(
                primary=RGBColor(0, 188, 212),      # 시안
                secondary=RGBColor(66, 66, 66),
                accent=RGBColor(255, 64, 129),      # 핑크
                background=RGBColor(33, 33, 33),    # 다크
                text_dark=RGBColor(255, 255, 255),
                text_light=RGBColor(255, 255, 255),
                gradient_start=RGBColor(33, 33, 33),
                gradient_end=RGBColor(66, 66, 66),
            ),
            "minimal": ColorTheme(
                primary=RGBColor(51, 51, 51),
                secondary=RGBColor(128, 128, 128),
                accent=RGBColor(0, 150, 136),       # 틸
                background=RGBColor(250, 250, 250),
                text_dark=RGBColor(51, 51, 51),
                text_light=RGBColor(255, 255, 255),
                gradient_start=RGBColor(240, 240, 240),
                gradient_end=RGBColor(255, 255, 255),
            ),
        }
        return themes.get(theme_name, themes["modern_blue"])

    def design(self, content: PresentationContent) -> Presentation:
        """
        요약된 콘텐츠를 PPT로 변환

        Args:
            content: 요약된 프레젠테이션 콘텐츠

        Returns:
            디자인된 Presentation 객체
        """
        self.prs = Presentation()

        # 슬라이드 크기 설정 (16:9)
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

        # 표지 슬라이드
        self._add_title_slide(content.title, content.subtitle)

        # 목차 슬라이드
        if len(content.slides) > 3:
            self._add_toc_slide(content.slides)

        # 콘텐츠 슬라이드
        for slide_content in content.slides:
            self._add_content_slide(slide_content)

        # 마무리 슬라이드
        self._add_closing_slide()

        return self.prs

    def _add_title_slide(self, title: str, subtitle: Optional[str] = None):
        """표지 슬라이드 생성"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 배경 도형 (상단 색상 바)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme.primary
        shape.line.fill.background()

        # 하단 악센트 라인
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(3),
            self.SLIDE_WIDTH, Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.accent
        line.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(3.5),
            self.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title or "프레젠테이션"
        p.font.name = self.FONT_TITLE
        p.font.size = self.FONT_SIZE_MAIN_TITLE
        p.font.bold = True
        p.font.color.rgb = self.theme.text_dark
        p.alignment = PP_ALIGN.LEFT

        # 부제목
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.MARGIN_LEFT, Inches(5.2),
                self.CONTENT_WIDTH, Inches(0.5)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle or ""
            p.font.name = self.FONT_BODY
            p.font.size = self.FONT_SIZE_SUBTITLE
            p.font.color.rgb = self.theme.secondary
            p.alignment = PP_ALIGN.LEFT

        # 날짜 (하단)
        from datetime import datetime
        date_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(6.5),
            Inches(3), Inches(0.4)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y년 %m월")
        p.font.name = self.FONT_BODY
        p.font.size = self.FONT_SIZE_SMALL
        p.font.color.rgb = self.theme.secondary

    def _add_toc_slide(self, slides: List[SlideContent]):
        """목차 슬라이드 생성"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 상단 바
        self._add_header_bar(slide)

        # 제목
        self._add_slide_title(slide, "목차", "CONTENTS")

        # 목차 항목들 (주요 섹션만)
        unique_titles = []
        for s in slides:
            # 중복 제거 및 번호 제거
            slide_title = s.title or "슬라이드"
            base_title = slide_title.split(" (")[0].split(" - ")[0]
            if base_title not in unique_titles:
                unique_titles.append(base_title)

        # 최대 8개
        unique_titles = unique_titles[:8]

        # 2열 레이아웃
        col_width = Inches(5.5)
        items_per_col = (len(unique_titles) + 1) // 2

        for i, title in enumerate(unique_titles):
            col = i // items_per_col
            row = i % items_per_col

            x = self.MARGIN_LEFT + (col * Inches(6))
            y = self.CONTENT_TOP + Inches(0.3) + (row * Inches(0.7))

            # 번호 원
            num_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y,
                Inches(0.35), Inches(0.35)
            )
            num_shape.fill.solid()
            num_shape.fill.fore_color.rgb = self.theme.primary
            num_shape.line.fill.background()

            # 번호 텍스트
            num_shape.text_frame.paragraphs[0].text = str(i + 1)
            num_shape.text_frame.paragraphs[0].font.color.rgb = self.theme.text_light
            num_shape.text_frame.paragraphs[0].font.size = Pt(12)
            num_shape.text_frame.paragraphs[0].font.bold = True
            num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 제목 텍스트
            title_box = slide.shapes.add_textbox(
                x + Inches(0.5), y,
                col_width, Inches(0.4)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title or "슬라이드"
            p.font.name = self.FONT_BODY
            p.font.size = self.FONT_SIZE_BODY
            p.font.color.rgb = self.theme.text_dark

    def _add_content_slide(self, content: SlideContent):
        """콘텐츠 슬라이드 생성"""
        layout_hint = content.layout_hint

        if layout_hint == "title_only":
            self._add_section_slide(content)
        elif layout_hint == "metrics":
            self._add_metrics_slide(content)
        elif layout_hint == "table":
            self._add_table_slide(content)
        else:
            self._add_bullet_slide(content)

    def _add_section_slide(self, content: SlideContent):
        """섹션 구분 슬라이드"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.primary
        bg.line.fill.background()

        # 제목 (가운데)
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(3),
            self.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = content.title or "섹션"
        p.font.name = self.FONT_TITLE
        p.font.size = self.FONT_SIZE_MAIN_TITLE
        p.font.bold = True
        p.font.color.rgb = self.theme.text_light
        p.alignment = PP_ALIGN.CENTER

    def _add_bullet_slide(self, content: SlideContent):
        """불릿 포인트 슬라이드"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 상단 바
        self._add_header_bar(slide)

        # 제목
        self._add_slide_title(slide, content.title)

        # 불릿 포인트들
        if content.bullets:
            y_offset = self.CONTENT_TOP + Inches(0.2)
            bullet_height = Inches(0.8)

            for i, bullet in enumerate(content.bullets[:self.MAX_LINES_PER_BULLET * 3]):
                # 불릿 아이콘 (작은 사각형)
                icon = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    self.MARGIN_LEFT, y_offset + Inches(0.08),
                    Inches(0.15), Inches(0.15)
                )
                icon.fill.solid()
                icon.fill.fore_color.rgb = self.theme.accent
                icon.line.fill.background()

                # 텍스트
                text_box = slide.shapes.add_textbox(
                    self.MARGIN_LEFT + Inches(0.35), y_offset,
                    self.CONTENT_WIDTH - Inches(0.5), bullet_height
                )
                tf = text_box.text_frame
                tf.word_wrap = True

                # 텍스트 줄 수 제한
                display_text = self._limit_text_lines(bullet, self.MAX_CHARS_PER_LINE, self.MAX_LINES_PER_BULLET)

                p = tf.paragraphs[0]
                p.text = display_text or ""
                p.font.name = self.FONT_BODY
                p.font.size = self.FONT_SIZE_BULLET
                p.font.color.rgb = self.theme.text_dark
                p.line_spacing = 1.2

                y_offset += bullet_height

                # 페이지 초과 방지
                if y_offset > Inches(6.5):
                    break

    def _add_metrics_slide(self, content: SlideContent):
        """메트릭/KPI 슬라이드"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 상단 바
        self._add_header_bar(slide)

        # 제목
        self._add_slide_title(slide, content.title)

        # 메트릭 카드들
        metrics = content.key_metrics[:4]
        num_metrics = len(metrics)

        if num_metrics == 0:
            return

        card_width = Inches(2.8)
        card_height = Inches(2)
        total_width = card_width * num_metrics + Inches(0.3) * (num_metrics - 1)
        start_x = (self.SLIDE_WIDTH - total_width) / 2

        for i, metric in enumerate(metrics):
            x = start_x + i * (card_width + Inches(0.3))
            y = self.CONTENT_TOP + Inches(0.5)

            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y,
                card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(245, 245, 245)
            card.line.color.rgb = RGBColor(230, 230, 230)

            # 상단 악센트 바
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y,
                card_width, Inches(0.08)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = self.theme.primary
            accent_bar.line.fill.background()

            # 값
            value_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.4),
                card_width - Inches(0.4), Inches(0.8)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric.get("value") or ""
            p.font.name = self.FONT_ENGLISH
            p.font.size = self.FONT_SIZE_METRIC_VALUE
            p.font.bold = True
            p.font.color.rgb = self.theme.primary
            p.alignment = PP_ALIGN.CENTER

            # 레이블
            label_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(1.3),
                card_width - Inches(0.4), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = metric.get("label") or ""
            p.font.name = self.FONT_BODY
            p.font.size = self.FONT_SIZE_METRIC_LABEL
            p.font.color.rgb = self.theme.secondary
            p.alignment = PP_ALIGN.CENTER

        # 추가 불릿 (있는 경우)
        if content.bullets:
            y_offset = self.CONTENT_TOP + Inches(3)
            for bullet in content.bullets[:3]:
                text_box = slide.shapes.add_textbox(
                    self.MARGIN_LEFT, y_offset,
                    self.CONTENT_WIDTH, Inches(0.5)
                )
                tf = text_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {bullet or ''}"
                p.font.name = self.FONT_BODY
                p.font.size = self.FONT_SIZE_SMALL
                p.font.color.rgb = self.theme.text_dark
                y_offset += Inches(0.5)

    def _add_table_slide(self, content: SlideContent):
        """표 슬라이드"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 상단 바
        self._add_header_bar(slide)

        # 제목
        self._add_slide_title(slide, content.title)

        table_data = content.table_data
        if not table_data:
            return

        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        if not headers:
            return

        # 표 크기 계산
        num_cols = len(headers)
        num_rows = min(len(rows) + 1, 7)  # 헤더 포함 최대 7행

        col_width = self.CONTENT_WIDTH / num_cols
        row_height = Inches(0.5)

        # 표 생성
        table = slide.shapes.add_table(
            num_rows, num_cols,
            self.MARGIN_LEFT, self.CONTENT_TOP + Inches(0.2),
            self.CONTENT_WIDTH, row_height * num_rows
        ).table

        # 열 너비 설정
        for i in range(num_cols):
            table.columns[i].width = int(col_width)

        # 헤더 행
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = self._truncate_text(header or "", 20)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.primary

            para = cell.text_frame.paragraphs[0]
            para.font.name = self.FONT_BODY
            para.font.size = self.FONT_SIZE_SMALL
            para.font.bold = True
            para.font.color.rgb = self.theme.text_light
            para.alignment = PP_ALIGN.CENTER

        # 데이터 행
        for row_idx, row_data in enumerate(rows[:num_rows - 1]):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    break
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = self._truncate_text(str(cell_text) if cell_text is not None else "", 25)

                # 교대 배경색
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)

                para = cell.text_frame.paragraphs[0]
                para.font.name = self.FONT_BODY
                para.font.size = self.FONT_SIZE_SMALL
                para.font.color.rgb = self.theme.text_dark
                para.alignment = PP_ALIGN.CENTER

    def _add_closing_slide(self):
        """마무리 슬라이드"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.primary
        bg.line.fill.background()

        # 감사 메시지
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(3),
            self.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "감사합니다"
        p.font.name = self.FONT_TITLE
        p.font.size = self.FONT_SIZE_MAIN_TITLE
        p.font.bold = True
        p.font.color.rgb = self.theme.text_light
        p.alignment = PP_ALIGN.CENTER

        # 부제
        subtitle_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT, Inches(4.5),
            self.CONTENT_WIDTH, Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Q & A"
        p.font.name = self.FONT_ENGLISH
        p.font.size = self.FONT_SIZE_SUBTITLE
        p.font.color.rgb = self.theme.text_light
        p.alignment = PP_ALIGN.CENTER

    def _add_header_bar(self, slide):
        """상단 헤더 바 추가"""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.primary
        bar.line.fill.background()

    def _add_slide_title(self, slide, title: str, subtitle: Optional[str] = None):
        """슬라이드 제목 추가"""
        # 좌측 악센트 바
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            self.MARGIN_LEFT, Inches(0.5),
            Inches(0.08), Inches(0.8)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.theme.primary
        accent.line.fill.background()

        # 제목
        title_box = slide.shapes.add_textbox(
            self.MARGIN_LEFT + Inches(0.25), Inches(0.5),
            self.CONTENT_WIDTH, Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title or "슬라이드"
        p.font.name = self.FONT_TITLE
        p.font.size = self.FONT_SIZE_SLIDE_TITLE
        p.font.bold = True
        p.font.color.rgb = self.theme.text_dark

        # 부제목
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.MARGIN_LEFT + Inches(0.25), Inches(1.1),
                self.CONTENT_WIDTH, Inches(0.4)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle or ""
            p.font.name = self.FONT_ENGLISH
            p.font.size = self.FONT_SIZE_SMALL
            p.font.color.rgb = self.theme.secondary

    def _limit_text_lines(self, text: str, chars_per_line: int, max_lines: int) -> str:
        """텍스트 줄 수 제한"""
        if not text:
            return ""
        words = text.split()
        lines = []
        current_line = ""

        for word in words:
            if len(current_line) + len(word) + 1 <= chars_per_line:
                current_line += (" " + word) if current_line else word
            else:
                if current_line:
                    lines.append(current_line)
                current_line = word

                if len(lines) >= max_lines:
                    break

        if current_line and len(lines) < max_lines:
            lines.append(current_line)

        result = "\n".join(lines)

        # 최대 줄 수 초과시 말줄임
        if len(lines) >= max_lines and len(words) > len(" ".join(lines).split()):
            result = result.rstrip() + "..."

        return result

    def _truncate_text(self, text: str, max_length: int) -> str:
        """텍스트 길이 제한"""
        if not text:
            return ""
        if len(text) <= max_length:
            return text
        return text[:max_length - 3] + "..."


def design_presentation(content: PresentationContent, theme: str = "modern_blue") -> Presentation:
    """
    편의 함수: 콘텐츠를 PPT로 디자인
    """
    designer = PPTDesigner(theme=theme)
    return designer.design(content)
